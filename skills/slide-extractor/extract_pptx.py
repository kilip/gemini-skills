# /// script
# dependencies = [
#   "python-pptx",
#   "pytesseract",
#   "easyocr",
#   "Pillow",
#   "numpy",
#   "opencv-python-headless",
# ]
# ///

import sys
import os
import io
import re
import numpy as np
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

# Global OCR Engines
TESSERACT_AVAILABLE = False
EASYOCR_READER = None

import pytesseract
import shutil
import platform

# Try to find tesseract in PATH first
tess_path = shutil.which("tesseract")

# Windows specific fallbacks if not in PATH
if not tess_path and platform.system() == "Windows":
    COMMON_TESS_PATHS = [
        r"C:\Program Files\Tesseract-OCR\tesseract.exe",
        r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
        os.path.expanduser(r"~\AppData\Local\Tesseract-OCR\tesseract.exe"),
    ]
    for path in COMMON_TESS_PATHS:
        if os.path.exists(path):
            tess_path = path
            break

if tess_path:
    pytesseract.pytesseract.tesseract_cmd = tess_path

try:
    pytesseract.get_tesseract_version()
    TESSERACT_AVAILABLE = True
except Exception:
    TESSERACT_AVAILABLE = False


def get_easyocr():
    global EASYOCR_READER
    if EASYOCR_READER is None:
        import easyocr
        EASYOCR_READER = easyocr.Reader(['en', 'id'], gpu=False)
    return EASYOCR_READER


def extract_text_from_pptx(pptx_path, preferred_engine=None):
    if not os.path.exists(pptx_path):
        print(f"Error: File not found at {pptx_path}", file=sys.stderr)
        return None

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"Error opening PPTX: {e}", file=sys.stderr)
        return None

    def perform_ocr(img_blob):
        # Priority 1: Tesseract
        if (preferred_engine is None or preferred_engine == "tesseract") and TESSERACT_AVAILABLE:
            try:
                img = Image.open(io.BytesIO(img_blob))
                ocr_text = pytesseract.image_to_string(img).strip()
                if ocr_text:
                    return f"[OCR-Tesseract]\n{ocr_text}"
            except Exception:
                pass

        # Priority 2: EasyOCR
        if preferred_engine is None or preferred_engine == "easyocr":
            try:
                reader = get_easyocr()
                img = Image.open(io.BytesIO(img_blob)).convert("RGB")
                img_array = np.array(img)
                results = reader.readtext(img_array, detail=0)
                ocr_text = " ".join(results).strip()
                if ocr_text:
                    return f"[OCR-EasyOCR]\n{ocr_text}"
            except Exception:
                pass

        return None

    def get_text_and_images(shape, text_runs):
        if hasattr(shape, "text") and shape.text.strip():
            text_runs.append(shape.text.strip())

        if shape.has_table:
            for row in shape.table.rows:
                row_text = [cell.text_frame.text.strip() for cell in row.cells if cell.text_frame.text.strip()]
                if row_text:
                    text_runs.append(" | ".join(row_text))

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for s in shape.shapes:
                get_text_and_images(s, text_runs)

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            ocr_res = perform_ocr(shape.image.blob)
            if ocr_res:
                text_runs.append(ocr_res)

    # Build output as a list of lines
    lines = []

    # Metadata header
    base_name = os.path.splitext(os.path.basename(pptx_path))[0]
    lines.append("---")
    lines.append(f"title: {base_name}")
    lines.append(f"filename: {os.path.basename(pptx_path)}")
    lines.append(f"slide_count: {len(prs.slides)}")
    lines.append("---")
    lines.append("")

    for i, slide in enumerate(prs.slides):
        lines.append(f"# Slide {i + 1}")
        text_runs = []
        for shape in slide.shapes:
            get_text_and_images(shape, text_runs)

        if text_runs:
            unique_text = []
            seen = set()
            for t in text_runs:
                clean_t = t.strip().lower()
                if clean_t and clean_t not in seen:
                    unique_text.append(t)
                    seen.add(clean_t)
            lines.append("\n".join(unique_text))
        else:
            lines.append("*[Empty Slide]*")

        lines.append("")
        lines.append("---")
        lines.append("")

    return "\n".join(lines)


if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser(description="Extract text from PPTX. Output saved as .md in the same folder.")
    parser.add_argument("pptx_path", help="Path to the .pptx file (must be a local file path)")
    parser.add_argument("--engine", choices=["tesseract", "easyocr"], help="Force specific OCR engine")
    args = parser.parse_args()

    pptx_path = args.pptx_path

    if not os.path.exists(pptx_path):
        print(f"Error: File not found: {pptx_path}", file=sys.stderr)
        sys.exit(1)

    if not pptx_path.lower().endswith(".pptx"):
        print(f"Error: File must be a .pptx file: {pptx_path}", file=sys.stderr)
        sys.exit(1)

    # Determine output path: same directory, same name, .md extension
    output_path = os.path.splitext(pptx_path)[0] + ".md"

    print(f"Extracting: {pptx_path}", file=sys.stderr)
    print(f"Output:     {output_path}", file=sys.stderr)

    result = extract_text_from_pptx(pptx_path, args.engine)

    if result is None:
        print("Extraction failed.", file=sys.stderr)
        sys.exit(1)

    with open(output_path, "w", encoding="utf-8") as f:
        f.write(result)

    print(f"Done. Read the output file: {output_path}", file=sys.stderr)